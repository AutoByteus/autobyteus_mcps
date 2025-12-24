import os
import pytest
from pptx import Presentation
from PIL import Image as PILImage
from app.main import create_ppt_from_images, replace_slide_with_image, append_images_as_slides, combine_images

@pytest.fixture
def test_images():
    return ["tests/test_red.png", "tests/test_blue.png"]

@pytest.fixture
def output_pptx(tmp_path):
    return str(tmp_path / "test_output.pptx")

@pytest.fixture
def output_image(tmp_path):
    return str(tmp_path / "combined_output.png")

def test_create_ppt_from_images(test_images, output_pptx):
    result = create_ppt_from_images(test_images, output_pptx)
    assert "output_path" in result
    assert os.path.exists(output_pptx)
    
    prs = Presentation(output_pptx)
    assert len(prs.slides) == 2

def test_replace_slide_with_image(test_images, output_pptx):
    # Create initial PPT
    create_ppt_from_images([test_images[0]], output_pptx)
    
    # Replace slide 0 with blue image
    replace_result = replace_slide_with_image(output_pptx, 0, test_images[1])
    assert "output_path" in replace_result
    
    # Verify content (simplified check)
    prs = Presentation(replace_result["output_path"])
    assert len(prs.slides) == 1

def test_append_images_as_slides(test_images, output_pptx):
    # Create initial PPT
    create_ppt_from_images([test_images[0]], output_pptx)
    
    # Append another image
    append_result = append_images_as_slides(output_pptx, [test_images[1]])
    assert append_result["appended"] == 1
    
    prs = Presentation(append_result["output_path"])
    assert len(prs.slides) == 2

def test_combine_images(test_images, output_image):
    # Test Vertical (Default)
    result = combine_images(test_images, output_image)
    assert "output_path" in result
    assert os.path.exists(output_image)
    assert result["image_count"] == 2
    assert result["direction"] == "vertical"
    
    with PILImage.open(output_image) as img:
        assert img.width > 0
        assert img.height > 0
        vert_height = img.height

    # Test Horizontal
    os.remove(output_image)
    result_horiz = combine_images(test_images, output_image, direction="horizontal")
    assert result_horiz["direction"] == "horizontal"
    
    with PILImage.open(output_image) as img:
        assert img.width > 0
        assert img.height > 0
        # Horizontal width should typically be larger than vertical width for the same set of similar images
        # But specifically, height should be less than vertical stacked height
        assert img.height < vert_height

