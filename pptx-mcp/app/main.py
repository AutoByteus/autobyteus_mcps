import logging
import os
from dataclasses import dataclass
from typing import Iterable

from mcp.server.fastmcp import FastMCP
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches

# Configure basic logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# --- Workspace Handling ---
# If AUTOBYTEUS_AGENT_WORKSPACE is set, change CWD to that path.

def initialize_workspace() -> None:
    workspace_path = os.environ.get('AUTOBYTEUS_AGENT_WORKSPACE')
    if workspace_path:
        logger.info("AUTOBYTEUS_AGENT_WORKSPACE found: '%s'", workspace_path)
        if os.path.isdir(workspace_path):
            try:
                os.chdir(workspace_path)
                logger.info("Successfully changed CWD to: %s", os.getcwd())
            except Exception as exc:
                logger.error("Failed to change CWD to '%s': %s", workspace_path, exc, exc_info=True)
        else:
            logger.warning("Workspace path '%s' does not exist or is not a directory. CWD not changed.", workspace_path)
    else:
        logger.info("AUTOBYTEUS_AGENT_WORKSPACE not set. Using default CWD.")


initialize_workspace()
# --- End Workspace Handling ---

mcp = FastMCP("PPTXImageServer")


@dataclass(frozen=True)
class SlideSize:
    width_emu: int
    height_emu: int
    ratio: float


def _get_image_size(image_path: str) -> tuple[int, int]:
    with PILImage.open(image_path) as image:
        return image.size


def _compute_slide_size_from_first_image(image_path: str, base_height_inches: float) -> SlideSize:
    width_px, height_px = _get_image_size(image_path)
    if height_px == 0:
        raise ValueError(f"Invalid image height for '{image_path}'.")
    ratio = width_px / height_px
    height_emu = int(Inches(base_height_inches))
    width_emu = int(height_emu * ratio)
    return SlideSize(width_emu=width_emu, height_emu=height_emu, ratio=ratio)


def _scale_to_contain(slide_size: SlideSize, image_width_px: int, image_height_px: int) -> tuple[int, int]:
    if image_width_px == 0 or image_height_px == 0:
        raise ValueError("Image dimensions must be non-zero.")
    scale = min(slide_size.width_emu / image_width_px, slide_size.height_emu / image_height_px)
    return int(image_width_px * scale), int(image_height_px * scale)


def _center_offsets(slide_size: SlideSize, shape_width: int, shape_height: int) -> tuple[int, int]:
    left = int((slide_size.width_emu - shape_width) / 2)
    top = int((slide_size.height_emu - shape_height) / 2)
    return left, top


def _derive_output_path(input_path: str, suffix: str = "_out") -> str:
    base, ext = os.path.splitext(input_path)
    return f"{base}{suffix}{ext or '.pptx'}"


def _clear_slide(slide) -> None:
    # Remove all shapes from slide to replace with a single image.
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)


def _validate_images(images: Iterable[str]) -> list[str]:
    validated = []
    for image_path in images:
        if not os.path.isfile(image_path):
            raise FileNotFoundError(f"Image not found: {image_path}")
        validated.append(image_path)
    return validated


def _insert_image_slide(prs: Presentation, slide_size: SlideSize, image_path: str) -> None:
    image_width_px, image_height_px = _get_image_size(image_path)
    shape_width, shape_height = _scale_to_contain(slide_size, image_width_px, image_height_px)
    left, top = _center_offsets(slide_size, shape_width, shape_height)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(image_path, left, top, width=shape_width, height=shape_height)


@mcp.tool(
    name="create_ppt_from_images",
    description="Create a new PowerPoint presentation from a list of images. Each image becomes a new slide. The slide dimensions are set based on the first image's aspect ratio. All images are scaled to fit within the slide boundaries without cropping."
)
def create_ppt_from_images(images: list[str], output_path: str) -> dict:
    """
    Create a new PPTX from a list of images.
    Slide size is determined by the first image's aspect ratio.
    Images are scaled to fit within the slide (no cropping).
    """
    base_height_inches = 7.5
    logger.info("Creating PPTX from %d images.", len(images))
    if not images:
        return {"error": "images list cannot be empty"}

    images = _validate_images(images)
    slide_size = _compute_slide_size_from_first_image(images[0], base_height_inches)

    prs = Presentation()
    prs.slide_width = slide_size.width_emu
    prs.slide_height = slide_size.height_emu

    for image_path in images:
        _insert_image_slide(prs, slide_size, image_path)

    prs.save(output_path)
    logger.info("Saved PPTX to %s", output_path)
    return {
        "output_path": output_path,
        "slide_width_emu": slide_size.width_emu,
        "slide_height_emu": slide_size.height_emu,
        "image_count": len(images),
    }


@mcp.tool(
    name="replace_slide_with_image",
    description="Replace the content of a specific slide (by 0-based index) with a single image. The image is scaled to fit within the existing slide dimensions without cropping. The original slide content is cleared."
)
def replace_slide_with_image(pptx_path: str, slide_index: int, image_path: str, output_path: str | None = None) -> dict:
    """
    Replace a slide (0-based index) with a single image.
    The image is scaled to fit within the slide (no cropping).
    """
    logger.info("Replacing slide %d in %s", slide_index, pptx_path)
    if not os.path.isfile(pptx_path):
        return {"error": f"pptx file not found: {pptx_path}"}
    _validate_images([image_path])

    prs = Presentation(pptx_path)
    if slide_index < 0 or slide_index >= len(prs.slides):
        return {"error": f"slide_index out of range: {slide_index}"}

    slide_size = SlideSize(width_emu=prs.slide_width, height_emu=prs.slide_height, ratio=prs.slide_width / prs.slide_height)
    slide = prs.slides[slide_index]
    _clear_slide(slide)

    image_width_px, image_height_px = _get_image_size(image_path)
    shape_width, shape_height = _scale_to_contain(slide_size, image_width_px, image_height_px)
    left, top = _center_offsets(slide_size, shape_width, shape_height)
    slide.shapes.add_picture(image_path, left, top, width=shape_width, height=shape_height)

    output_path = output_path or _derive_output_path(pptx_path)
    prs.save(output_path)
    logger.info("Saved PPTX to %s", output_path)
    return {
        "output_path": output_path,
        "slide_index": slide_index,
    }


@mcp.tool(
    name="append_images_as_slides",
    description="Append one or more images as new slides to the end of an existing PowerPoint presentation. Each image gets its own slide, scaled to fit the presentation's dimensions without cropping."
)
def append_images_as_slides(pptx_path: str, images: list[str], output_path: str | None = None) -> dict:
    """
    Append images as new slides in an existing PPTX.
    Images are scaled to fit within the slide (no cropping).
    """
    logger.info("Appending %d images to %s", len(images), pptx_path)
    if not images:
        return {"error": "images list cannot be empty"}
    if not os.path.isfile(pptx_path):
        return {"error": f"pptx file not found: {pptx_path}"}

    images = _validate_images(images)

    prs = Presentation(pptx_path)
    slide_size = SlideSize(width_emu=prs.slide_width, height_emu=prs.slide_height, ratio=prs.slide_width / prs.slide_height)

    for image_path in images:
        _insert_image_slide(prs, slide_size, image_path)

    output_path = output_path or _derive_output_path(pptx_path)
    prs.save(output_path)
    logger.info("Saved PPTX to %s", output_path)
    return {
        "output_path": output_path,
        "appended": len(images),
    }


@mcp.tool(
    name="combine_images",
    description="Combine multiple images into a single image. Supports 'vertical' (top-to-bottom) or 'horizontal' (left-to-right) arrangement. The output dimensions are calculated to fit all images."
)
def combine_images(images: list[str], output_path: str, direction: str = "vertical") -> dict:
    """
    Combine images into a single file vertically or horizontally.
    """
    logger.info("Combining %d images (direction=%s).", len(images), direction)
    if not images:
        return {"error": "images list cannot be empty"}
    
    if direction not in ("vertical", "horizontal"):
        return {"error": f"Invalid direction '{direction}'. Must be 'vertical' or 'horizontal'."}

    images = _validate_images(images)

    opened_images = []
    try:
        for img_path in images:
            opened_images.append(PILImage.open(img_path))
        
        if not opened_images:
             return {"error": "No images loaded"}

        # Calculate dimensions based on direction
        if direction == "vertical":
            max_width = max(img.width for img in opened_images)
            total_height = sum(img.height for img in opened_images)
            total_width = max_width
        else: # horizontal
            max_height = max(img.height for img in opened_images)
            total_width = sum(img.width for img in opened_images)
            total_height = max_height

        # Create new image
        new_im = PILImage.new('RGB', (total_width, total_height), (255, 255, 255))

        x_offset = 0
        y_offset = 0
        for img in opened_images:
            new_im.paste(img, (x_offset, y_offset))
            if direction == "vertical":
                y_offset += img.height
            else:
                x_offset += img.width

        new_im.save(output_path)
        logger.info("Saved combined image to %s", output_path)
        
        return {
            "output_path": output_path,
            "total_width": total_width,
            "total_height": total_height,
            "image_count": len(images),
            "direction": direction
        }
    except Exception as e:
        logger.error("Failed to combine images: %s", e)
        return {"error": str(e)}
    finally:
        # Close images to free resources
        for img in opened_images:
            try:
                img.close()
            except:
                pass


def run_server() -> None:
    logger.info("Starting PPTX Image MCP server with stdio transport...")
    mcp.run(transport="stdio")


if __name__ == "__main__":
    run_server()
